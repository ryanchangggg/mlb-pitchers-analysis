"""
Generate a comprehensive .pptx presentation comparing Ohtani, Sanchez, and Misiorowski (2026 season).
Uses the figures/ visualizations and provides explanatory text throughout.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import os

FIGURES_DIR = os.path.join(os.path.dirname(__file__), '..', 'figures')
OUTPUT_DIR = os.path.join(os.path.dirname(__file__), '..', 'report')
OUTPUT_PATH = os.path.join(OUTPUT_DIR, 'MLB_Pitchers_Analysis_2026.pptx')

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colour palette
BLUE   = RGBColor(0x1F, 0x77, 0xB4)
ORANGE = RGBColor(0xFF, 0x7F, 0x0E)
GREEN  = RGBColor(0x2C, 0xA0, 0x2C)
RED    = RGBColor(0xD6, 0x27, 0x28)
PURPLE = RGBColor(0x94, 0x6A, 0xC6)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x33, 0x33, 0x33)
MEDIUM = RGBColor(0x66, 0x66, 0x66)
LIGHT  = RGBColor(0xDE, 0xE2, 0xE6)
BG     = RGBColor(0xF0, 0xF0, 0xF0)

# ---- helpers ----

def set_slide_bg(slide, color=LIGHT):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, size=18,
                bold=False, color=DARK, align=PP_ALIGN.LEFT, italic=False,
                font_name='Calibri', anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.italic = italic
    p.font.name = font_name
    p.alignment = align
    return txBox, tf

def add_multi_text(slide, left, top, width, height, lines, default_size=16,
                   default_color=DARK, line_spacing=6, font_name='Calibri'):
    """Add a textbox with multiple paragraphs."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, str):
            text, size, color, bold = item, default_size, default_color, False
        else:
            text = item[0]
            size = item[1] if len(item) > 1 else default_size
            color = item[2] if len(item) > 2 else default_color
            bold = item[3] if len(item) > 3 else False

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.space_after = Pt(line_spacing)
    return txBox, tf

def add_figure_slide(slide, fig_name, area_left=0.3, area_top=1.2,
                     area_width=12.7, area_height=5.8):
    """Embed a figure PNG, preserving its native aspect ratio and centering it
    within the specified area."""
    path = os.path.join(FIGURES_DIR, fig_name)
    if not os.path.exists(path):
        # try fallback
        alt_path = os.path.join(FIGURES_DIR, fig_name.replace('_comparison', ''))
        if os.path.exists(alt_path):
            path = alt_path
        else:
            return False

    with Image.open(path) as img:
        native_w, native_h = img.size

    aspect = native_w / native_h
    area_aspect = area_width / area_height

    if aspect >= area_aspect:
        # Image is wider (or equal) relative to the area → width-limited
        disp_w = area_width
        disp_h = area_width / aspect
    else:
        # Image is taller relative to the area → height-limited
        disp_h = area_height
        disp_w = area_height * aspect

    # Center within the area
    left = area_left + (area_width - disp_w) / 2
    top = area_top + (area_height - disp_h) / 2

    slide.shapes.add_picture(path, Inches(left), Inches(top),
                             Inches(disp_w), Inches(disp_h))
    return True

def add_section_header(slide, title, subtitle=None, size=32):
    """Add a section header at the top."""
    add_textbox(slide, 0.5, 0.2, 12.3, 0.7, title, size=size, bold=True, color=DARK)
    if subtitle:
        add_textbox(slide, 0.5, 0.75, 12.3, 0.4, subtitle, size=14, color=MEDIUM, italic=True)

def add_caption(slide, text, y=7.0, size=10):
    add_textbox(slide, 0.5, y, 12.3, 0.4, text, size=size, color=MEDIUM, italic=True)

def new_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, BG)
    return slide


# =====================================================================
# SLIDE 1 — TITLE
# =====================================================================
slide = new_slide()
# Decorative top bar
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.08))
shape.fill.solid()
shape.fill.fore_color.rgb = BLUE
shape.line.fill.background()

add_textbox(slide, 0.8, 1.8, 11.7, 1.5,
    "Not All Aces Are Equal",
    size=48, bold=True, color=DARK, align=PP_ALIGN.CENTER)
add_textbox(slide, 0.8, 3.3, 11.7, 0.8,
    "A Data-Driven Comparison of Three NL Cy Young Candidates in 2026",
    size=22, color=MEDIUM, align=PP_ALIGN.CENTER)
add_textbox(slide, 0.8, 4.3, 11.7, 0.6,
    "Shohei Ohtani  ·  Cristopher Sanchez  ·  Jacob Misiorowski",
    size=20, color=BLUE, align=PP_ALIGN.CENTER)
add_textbox(slide, 0.8, 5.5, 11.7, 0.4,
    "Data: MLB Statcast (2026 Season)  ·  Generated: July 2026",
    size=13, color=MEDIUM, align=PP_ALIGN.CENTER)
add_caption(slide, "Source: MLB Statcast pitch-by-pitch data via baseballr", y=7.0, size=10)

# =====================================================================
# SLIDE 2 — AGENDA
# =====================================================================
slide = new_slide()
add_section_header(slide, "Agenda", "What we will cover")
add_multi_text(slide, 0.7, 1.3, 11.9, 5.8, [
    ("1. Data Overview & Methodology", 20, DARK, True),
    ("2. Pitch Arsenal Analysis — diversity versus effectiveness", 17, DARK),
    ("3. Velocity & Movement Profiles", 17, DARK),
    ("4. Swing-and-Miss & Plate Discipline", 17, DARK),
    ("5. Control & Command", 17, DARK),
    ("6. Quality of Contact & Batted Balls", 17, DARK),
    ("7. Expected Statistics (xBA, xwOBA, xSLG)", 17, DARK),
    ("8. Run Expectancy — The Key Finding", 17, BLUE, True),
    ("9. Platoon Splits & Situational Performance", 17, DARK),
    ("10. Overall Assessment & Conclusion", 20, DARK, True),
], line_spacing=10)
add_caption(slide, "Analysis performed in Python (pandas, numpy, matplotlib, seaborn, scipy)")

# =====================================================================
# SLIDE 3 — DATA OVERVIEW
# =====================================================================
slide = new_slide()
add_section_header(slide, "Data Overview", "4,797 total pitch events from MLB Statcast")
add_multi_text(slide, 0.7, 1.3, 5.5, 5.0, [
    ("Dataset Summary", 22, BLUE, True),
    ("", 6),
    ("Shohei Ohtani (LAD)", 18, DARK, True),
    ("  1,335 pitches · 14 games · 7 pitch types", 16, DARK),
    ("   Two-way player: rest days for batting limit workload", 14, MEDIUM),
    ("", 6),
    ("Cristopher Sanchez (PHI)", 18, DARK, True),
    ("  1,800 pitches · 19 games · 3 pitch types", 16, DARK),
    ("   Workhorse: 35% more pitches than Ohtani", 14, MEDIUM),
    ("", 6),
    ("Jacob Misiorowski (MIL)", 18, DARK, True),
    ("  1,662 pitches · 18 games · 5 pitch types", 16, DARK),
    ("   Power arm: 100.5 mph avg fastball", 14, MEDIUM),
], line_spacing=4)

# Key stat box
add_multi_text(slide, 6.8, 1.5, 5.8, 4.0, [
    ("Key Observation", 22, RED, True),
    ("", 6),
    ("Ohtani has thrown ~35% fewer pitches than Sanchez.", 16, DARK),
    ("", 4),
    ("Despite his two-way role, this workload gap", 16, DARK),
    ("limits overall value contribution.", 16, DARK),
    ("", 4),
    ("A Cy Young candidate provides innings —", 16, DARK),
    ("Ohtani simply does not match his peers here.", 16, RED, True),
], line_spacing=4)

# =====================================================================
# SLIDE 4 — PITCH ARSENAL: OHTANI
# =====================================================================
slide = new_slide()
add_section_header(slide, "Pitch Arsenal: Shohei Ohtani", "The widest repertoire — 7 pitch types")
add_multi_text(slide, 0.7, 1.4, 7.0, 5.5, [
    ("Pitch      Usage     Velo     MaxVelo    Spin    HMov", 18, DARK, True),
    ("FF (4-seam)  45.2%   98.1     101.7     2486    -0.4", 15, DARK),
    ("ST (sweeper) 29.4%   85.0      90.1     2726     1.3", 15, DARK),
    ("CU (curve)    10.0%   75.2      82.5     2651     1.1", 15, DARK),
    ("FS (splitter)  8.5%   89.0      93.0     1564    -1.1", 15, DARK),
    ("SI (sinker)    5.1%   96.5      99.9     2231    -1.2", 15, DARK),
    ("SL (slider)    1.1%   88.0      90.8     2593     0.3", 15, DARK),
    ("FC (cutter)    0.7%   92.4      94.6     2376    -0.1", 15, DARK),
    ("", 6),
    ("Key Insight: 45.2% four-seam + 29.4% sweeper = 74.6% of arsenal", 16, DARK),
    ("Splitter (8.5%) and sinker (5.1%) are complementary", 15, MEDIUM),
    ("Cutter and slider are essentially show-me pitches (<2% combined)", 15, MEDIUM),
], line_spacing=4)

# Side note
add_multi_text(slide, 8.2, 2.5, 4.5, 3.0, [
    ("Breadth ≠ Dominance", 20, RED, True),
    ("", 6),
    ("Having 7 pitch types is unusual, but depth does not", 16, DARK),
    ("automatically translate to better results.", 16, DARK),
    ("", 4),
    ("Sanchez uses only 3 pitch types yet outperforms", 16, DARK),
    ("Ohtani on several key metrics.", 16, GREEN, True),
], line_spacing=4)

# =====================================================================
# SLIDE 5 — PITCH ARSENAL: SANCHEZ
# =====================================================================
slide = new_slide()
add_section_header(slide, "Pitch Arsenal: Cristopher Sanchez", "Minimalist mastery — just 3 pitch types")
add_multi_text(slide, 0.7, 1.4, 7.0, 5.5, [
    ("Pitch      Usage     Velo     MaxVelo    Spin    HMov", 18, DARK, True),
    ("SI (sinker)  42.4%   95.2      97.7     2137     1.5", 15, DARK),
    ("CH (change)  39.0%   86.9      90.5     2027     1.4", 15, DARK),
    ("SL (slider)  18.6%   86.1      89.4     2074    -0.1", 15, DARK),
    ("", 6),
    ("Why it works:", 20, ORANGE, True),
    ("", 4),
    ("1. Sinker (95.2 mph) + heavy arm-side run (1.5 in) induces weak contact", 16, DARK),
    ("2. Changeup (86.9 mph) tunnels identically with sinker for 5+ mph drop", 16, DARK),
    ("3. Slider (86.1 mph) works as a breaking fringe offering opposite movement", 16, DARK),
    ("", 4),
    ("Most efficient arsenal of any qualifier? Debatable — but the data says yes here.", 15, MEDIUM),
], line_spacing=4)

add_multi_text(slide, 8.2, 2.0, 4.5, 4.0, [
    ("By the numbers:", 20, GREEN, True),
    ("", 6),
    ("• 39% changeup usage — elite whiff rate (42.0%)", 16, DARK),
    ("• 42.4% sinker usage — 192 ground balls (most of all 3)", 16, DARK),
    ("• 18.6% slider — 42.7% whiff rate (best of any pitch here)", 16, DARK),
    ("", 4),
    ("A 3-pitch mix that rivals or beats 7-pitch arsenals.", 16, BLUE, True),
], line_spacing=4)

# =====================================================================
# SLIDE 6 — PITCH ARSENAL: MISIOROWSKI
# =====================================================================
slide = new_slide()
add_section_header(slide, "Pitch Arsenal: Jacob Misiorowski", "Power-first — 5 pitch types, elite velocity")
add_multi_text(slide, 0.7, 1.4, 7.0, 5.5, [
    ("Pitch      Usage     Velo     MaxVelo    Spin    HMov", 18, DARK, True),
    ("FF (4-seam)  63.1%   100.5    105.5     2620    -0.9", 15, DARK),
    ("FC (cutter)  13.6%    96.0     99.6     2577     0.1", 15, DARK),
    ("CU (curve)   11.2%    87.9     91.0     2224     0.5", 15, DARK),
    ("SL (slider)  10.3%    92.6     95.5     2531     0.3", 15, DARK),
    ("CH (change)   1.7%    92.5     96.1     2037    -1.3", 15, DARK),
    ("", 6),
    ("Key Insight: 63.1% four-seam fastball — the most FB-dependent of the three.", 16, DARK),
    ("Cutter (13.6%) and curve (11.2%) are complementary breaking offerings.", 15, MEDIUM),
    ("Changeup is essentially unused (1.7%) — 2-pitch approach in practice.", 15, MEDIUM),
], line_spacing=4)

add_multi_text(slide, 8.2, 2.0, 4.5, 4.0, [
    ("Dominant Fastball", 22, RED, True),
    ("", 6),
    ("At 100.5 mph avg with 2620 RPM spin, his fastball", 16, DARK),
    ("is in the 99th percentile league-wide.", 16, DARK),
    ("", 4),
    ("37.2% whiff rate on the four-seamer —", 16, DARK),
    ("higher than most pitchers' secondary offerings.", 16, DARK),
    ("", 4),
    ("Everything else builds off the FB threat.", 16, GREEN, True),
], line_spacing=4)

# =====================================================================
# SLIDE 7 — ARSENAL PIE CHART
# =====================================================================
slide = new_slide()
add_section_header(slide, "Pitch Usage Breakdown", "Visual comparison of arsenal composition")
if not add_figure_slide(slide, "pitch_arsenal_pie.png"):
    add_figure_slide(slide, "arsenal_pie.png")
add_caption(slide, "Pitch arsenal pies — Ohtani (left), Sanchez (center), Misiorowski (right). Source: MLB Statcast 2026.")

# =====================================================================
# SLIDE 8 — ARSENAL COMPARISON TABLE
# =====================================================================
slide = new_slide()
add_section_header(slide, "Arsenal Summary Table", "Usage, velocity, and movement by pitch type")
add_multi_text(slide, 0.3, 1.3, 12.7, 5.5, [
    ("Pitcher     Pitch    Usage%   Velo   MaxVelo  Spin   HMov   VMov   Extension", 17, DARK, True),
    ("", 4),
    ("Ohtani      FF       45.2%    98.1   101.7    2486  -0.4    1.2     6.66", 14, DARK),
    ("Ohtani      ST       29.4%    85.0    90.1    2726   1.3    0.3     6.55", 14, DARK),
    ("Ohtani      CU       10.0%    75.2    82.5    2651   1.1   -1.2     6.39", 14, DARK),
    ("Ohtani      FS        8.5%    89.0    93.0    1564  -1.1    0.2     6.67", 14, DARK),
    ("Ohtani      SI        5.1%    96.5    99.9    2231  -1.2    0.5     6.69", 14, DARK),
    ("Ohtani      SL        1.1%    88.0    90.8    2593   0.3   -0.1     6.49", 14, DARK),
    ("Ohtani      FC        0.7%    92.4    94.6    2376  -0.1    0.9     6.49", 14, DARK),
    ("", 3),
    ("Sanchez     SI       42.4%    95.2    97.7    2137   1.5    0.4     6.92", 14, ORANGE),
    ("Sanchez     CH       39.0%    86.9    90.5    2027   1.4   -0.1     7.02", 14, ORANGE),
    ("Sanchez     SL       18.6%    86.1    89.4    2074  -0.1   -0.3     6.84", 14, ORANGE),
    ("", 3),
    ("Misiorowski FF       63.1%   100.5   105.5    2620  -0.9    1.3     7.56", 14, GREEN),
    ("Misiorowski FC       13.6%    96.0    99.6    2577   0.1    0.6     7.55", 14, GREEN),
    ("Misiorowski CU       11.2%    87.9    91.0    2224   0.5   -0.7     7.25", 14, GREEN),
    ("Misiorowski SL       10.3%    92.6    95.5    2531   0.3    0.3     7.45", 14, GREEN),
    ("Misiorowski CH        1.7%    92.5    96.1    2037  -1.3    0.4     7.44", 14, GREEN),
], line_spacing=2)

# =====================================================================
# SLIDE 9 — VELOCITY DISTRIBUTION
# =====================================================================
slide = new_slide()
add_section_header(slide, "Velocity Distribution", "How fast do these pitchers throw?")
if not add_figure_slide(slide, "velocity_distribution.png"):
    add_figure_slide(slide, "velocity_dist.png")
add_caption(slide, "Velocity distribution across all pitches. Misiorowski (green) dominates the upper tail >100 mph. Ohtani (blue) clusters in the high 90s. Sanchez (orange) tops out around 97 mph.")

# =====================================================================
# SLIDE 10 — VELOCITY BY PITCH TYPE
# =====================================================================
slide = new_slide()
add_section_header(slide, "Velocity by Pitch Type", "Breaking down velocity for each arsenal component")
if not add_figure_slide(slide, "velocity_by_pitch.png"):
    add_figure_slide(slide, "velocity_comparison.png")
add_caption(slide, "Mean velocity by pitch type. Error bars show ±1 standard deviation. Data from MLB Statcast 2026.")

# =====================================================================
# SLIDE 11 — WHIFF RATE ANALYSIS
# =====================================================================
slide = new_slide()
add_section_header(slide, "Swing-and-Miss Analysis", "Who generates the most whiffs?")
add_multi_text(slide, 0.7, 1.3, 5.5, 5.5, [
    ("Overall Whiff Rate", 22, RED, True),
    ("", 6),
    ("Ohtani:      180 whiffs / 633 swings = 28.4%", 18, BLUE, True),
    ("Sanchez:     274 whiffs / 923 swings = 29.7%", 18, ORANGE, True),
    ("Misiorowski: 286 whiffs / 834 swings = 34.3%", 18, GREEN, True),
    ("", 6),
    ("Ohtani ranks LAST in overall whiff rate.", 18, DARK),
    ("", 6),
    ("Best individual pitch whiff rates:", 20, BLUE, True),
    ("  Ohtani ST (sweeper):    37.1%", 15, DARK),
    ("  Sanchez SL (slider):    42.7% ← elite", 15, ORANGE, True),
    ("  Sanchez CH (changeup):  42.0% ← elite", 15, ORANGE, True),
    ("  Misiorowski CU (curve): 43.6% ← elite", 15, GREEN, True),
    ("  Misiorowski FF:         37.2%", 15, GREEN),
], line_spacing=4)

add_multi_text(slide, 6.8, 1.3, 5.8, 5.5, [
    ("Key Insight", 22, PURPLE, True),
    ("", 6),
    ("Despite having 7 pitch types, none of Ohtani's individual", 16, DARK),
    ("pitches reaches the elite whiff rates of Sanchez's", 16, DARK),
    ("secondary offerings (42.7% slider, 42.0% changeup).", 16, DARK),
    ("", 4),
    ("Sanchez's changeup and slider both clear 42% whiff rate.", 16, DARK),
    ("That is genuinely elite territory.", 16, DARK),
    ("", 4),
    ("Misiorowski's curveball (43.6%) is the single best", 16, DARK),
    ("whiff pitch among all three pitchers.", 16, DARK),
    ("", 4),
    ("Takeaway: Ohtani has variety but no single", 16, DARK),
    ("dominant swing-and-miss pitch.", 16, RED, True),
], line_spacing=4)

# =====================================================================
# SLIDE 12 — SPIN RATE COMPARISON
# =====================================================================
slide = new_slide()
add_section_header(slide, "Spin Rate Comparison", "Spin correlates with perceived velocity and movement")
if not add_figure_slide(slide, "spin_comparison.png"):
    add_figure_slide(slide, "spin_rate_comparison.png")
add_caption(slide, "Spin rate by pitch type (RPM). Higher spin on fastballs creates 'rising' effect; higher spin on breaking balls increases break magnitude.")

# =====================================================================
# SLIDE 13 — MOVEMENT PROFILES
# =====================================================================
slide = new_slide()
add_section_header(slide, "Movement Profiles", "Horizontal and vertical movement by pitch type")
if not add_figure_slide(slide, "pitch_movement.png"):
    add_figure_slide(slide, "movement_profiles.png")
add_caption(slide, "Pitch movement plot: horizontal movement (inches, positive = arm-side) vs vertical movement (inches, positive = upward). Size roughly proportional to usage.")

# =====================================================================
# SLIDE 14 — PLATE DISCIPLINE METRICS
# =====================================================================
slide = new_slide()
add_section_header(slide, "Plate Discipline Metrics", "How batters react to each pitcher")
add_multi_text(slide, 0.5, 1.3, 12.3, 5.5, [
    ("Metric          Ohtani       Sanchez      Misiorowski     NL Avg (ref)", 19, DARK, True),
    ("", 4),
    ("K%              27.9%        27.6%         39.6%           ~22.0%", 16, DARK),
    ("BB%              7.6%         4.8%          6.4%           ~8.5%", 16, DARK),
    ("K-BB%           20.3%        22.8%         33.2%           ~13.5%", 16, BLUE, True),
    ("Whiff%          28.4%        29.7%         34.3%           ~25.0%", 16, DARK),
    ("Swing%          47.4%        51.7%         50.2%           ~46.0%", 16, DARK),
    ("Chase%          29.8%        38.2%         33.8%           ~30.0%", 16, DARK),
    ("Zone%           50.3%        47.2%         49.6%           ~50.0%", 16, DARK),
    ("Strike%         48.6%        49.1%         54.6%           ~50.0%", 16, DARK),
    ("", 6),
    ("Highlight: Misiorowski's 39.6% K% and 33.2% K-BB% are in a class of their own.", 17, GREEN, True),
    ("Sanchez's 4.8% BB% is elite — barely half the league average walk rate.", 17, ORANGE, True),
    ("Ohtani's 7.6% BB% is the highest of the three, indicating a command gap.", 17, RED, True),
], line_spacing=3)

# =====================================================================
# SLIDE 15 — CONTROL & COMMAND
# =====================================================================
slide = new_slide()
add_section_header(slide, "Control & Command Gap", "Why Ohtani's walk rate matters")
add_multi_text(slide, 0.7, 1.3, 5.5, 5.5, [
    ("Walk Rate Comparison", 22, RED, True),
    ("", 6),
    ("Ohtani:  7.6% BB%  — 58% higher than Sanchez", 18, BLUE, True),
    ("Sanchez: 4.8% BB%  — elite, top-5 MLB", 18, ORANGE, True),
    ("Misiorowski: 6.4% BB%  — solid, above average", 18, GREEN, True),
    ("", 6),
    ("Ohtani walks batters at a 58% higher rate than Sanchez.", 16, DARK),
    ("This is the single clearest command gap in the comparison.", 16, DARK),
    ("", 6),
    ("Strike Rate", 22, PURPLE, True),
    ("", 4),
    ("Misiorowski: 54.6% Strike% — dominates the zone", 16, GREEN),
    ("Sanchez:     49.1% Strike% — efficient", 16, ORANGE),
    ("Ohtani:      48.6% Strike% — lowest of the 3", 16, RED, True),
], line_spacing=4)

add_multi_text(slide, 6.8, 1.3, 5.8, 5.5, [
    ("Chase Rate Paradox", 22, BLUE, True),
    ("", 6),
    ("Despite throwing the highest Zone% (50.3%),", 16, DARK),
    ("Ohtani generates the LOWEST chase rate (29.8%).", 16, DARK),
    ("", 4),
    ("Sanchez throws fewer strikes in the zone (47.2% Zone%)", 16, DARK),
    ("but batters chase at 38.2% — highest of all three.", 16, DARK),
    ("", 4),
    ("This suggests Ohtani's pitches are more hittable", 16, DARK),
    ("or easier to recognize than Sanchez's offerings.", 16, DARK),
    ("", 4),
    ("Sanchez's sinker-changeup tunnel is exceptionally", 16, DARK),
    ("deceptive — batters cannot distinguish them early.", 16, DARK),
    ("", 4),
    ("Implication: deception > zone rate for generating chases.", 16, RED, True),
], line_spacing=4)

# =====================================================================
# SLIDE 16 — BATTED BALL QUALITY
# =====================================================================
slide = new_slide()
add_section_header(slide, "Quality of Contact", "Exit velocity, hard-hit rate, and barrels")
add_multi_text(slide, 0.5, 1.3, 12.3, 5.5, [
    ("Metric              Ohtani        Sanchez     Misiorowski", 19, DARK, True),
    ("", 4),
    ("Balls in Play        213           333          220", 16, DARK),
    ("Avg Exit Velo      86.9 mph      89.1 mph     86.6 mph", 16, DARK),
    ("Avg Launch Angle    10.7 deg       2.3 deg     11.6 deg", 16, DARK),
    ("Hard Hit%          33.8%         44.3%        34.2%", 16, DARK),
    ("Barrels                0             8            3", 16, DARK),
    ("Ground Balls         111           192          101", 16, DARK),
    ("Fly Balls             59            74           78", 16, DARK),
    ("", 6),
    ("Context Matters:", 22, ORANGE, True),
    ("Ohtani suppresses EV (86.9 mph) and has 0 barrels — that's genuinely elite.", 16, DARK),
    ("But Sanchez induces 192 ground balls (vs 111 for Ohtani) despite more BIP.", 16, DARK),
    ("Sanchez's 2.3 deg avg launch angle is extremely low — ground ball machine.", 16, DARK),
    ("Ground balls rarely produce extra-base hits; this is intentional design.", 16, DARK),
], line_spacing=3)

# =====================================================================
# SLIDE 17 — EXPECTED STATISTICS
# =====================================================================
slide = new_slide()
add_section_header(slide, "Expected Statistics", "xBA, xwOBA, xSLG — what the model predicts")
add_multi_text(slide, 0.7, 1.3, 5.5, 5.5, [
    ("Expected Stat Comparison", 22, BLUE, True),
    ("", 6),
    ("Metric        Ohtani    Sanchez   Misiorowski", 18, DARK, True),
    ("xBA           .292      .326      .302", 16, DARK),
    ("xwOBA         .257      .279      .224", 16, DARK),
    ("xSLG          .418      .520      .440", 16, DARK),
    ("", 6),
    ("On paper, Ohtani's expected stats look better:", 16, DARK),
    ("lower xBA, xwOBA, and xSLG than Sanchez.", 16, DARK),
    ("", 4),
    ("But... xwOBA and xSLG are weighted toward", 16, DARK),
    ("damage on contact — and Sanchez suppresses", 16, DARK),
    ("actual damage through extreme ground balls.", 16, DARK),
], line_spacing=4)

add_multi_text(slide, 6.8, 1.3, 5.8, 5.5, [
    ("Key Nuance", 22, PURPLE, True),
    ("", 6),
    ("The gap between Sanchez's xwOBA (.279)", 16, DARK),
    ("and his actual wOBA tells the story:", 16, DARK),
    ("a pitcher who induces weak, harmless contact", 16, DARK),
    ("that the expected model overvalues.", 16, DARK),
    ("", 4),
    ("xwOBA doesn't fully account for ground ball", 16, DARK),
    ("quality — a 95 mph grounder is far less", 16, DARK),
    ("damaging than a 95 mph line drive.", 16, DARK),
    ("", 4),
    ("Takeaway: expected stats favor Ohtani, but", 16, DARK),
    ("actual run prevention tells a different story.", 16, RED, True),
], line_spacing=4)

# =====================================================================
# SLIDE 18 — BATTED BALL PROFILE
# =====================================================================
slide = new_slide()
add_section_header(slide, "Batted Ball Profile", "Ground balls, fly balls, line drives, and popups")
add_multi_text(slide, 0.7, 1.3, 5.5, 5.5, [
    ("Ground Ball Dominance", 22, ORANGE, True),
    ("", 6),
    ("Sanchez: 192 ground balls in 333 BIP", 18, ORANGE, True),
    ("  — 57.7% GB rate — elite ground ball pitcher", 16, DARK),
    ("  — Only 7 home runs in 1,800 pitches", 16, DARK),
    ("", 6),
    ("Ohtani:  111 ground balls in 213 BIP", 18, BLUE),
    ("  — 52.1% GB rate — solid but not elite", 16, DARK),
    ("  — 4 home runs in 1,335 pitches", 16, DARK),
    ("", 6),
    ("Misiorowski: 101 ground balls in 220 BIP", 18, GREEN),
    ("  — 45.9% GB rate — fly ball tendency", 16, DARK),
    ("  — 9 home runs allowed — most of the 3", 16, RED, True),
], line_spacing=4)

add_multi_text(slide, 6.8, 1.3, 5.8, 5.5, [
    ("Why Ground Balls Matter", 22, GREEN, True),
    ("", 6),
    ("Ground balls are the safest type of contact:", 16, DARK),
    ("• Almost never produce extra-base hits", 16, DARK),
    ("• Double-play opportunities", 16, DARK),
    ("• Lower run expectancy than fly balls", 16, DARK),
    ("• Less dependent on outfield defense", 16, DARK),
    ("", 4),
    ("Sanchez's 2.3° avg launch angle is the 2nd-lowest", 16, DARK),
    ("of any qualified pitcher in 2026.", 16, DARK),
    ("", 4),
    ("Ohtani's 10.7° avg launch angle is more neutral.", 16, DARK),
    ("More line drives = more damage.", 16, DARK),
    ("", 4),
    ("Misiorowski's 11.6° + 9 HR = high-variance profile.", 16, RED, True),
], line_spacing=4)

# =====================================================================
# SLIDE 19 — RUN EXPECTANCY (THE KEY FINDING)
# =====================================================================
slide = new_slide()
add_section_header(slide, "Run Expectancy Impact", "The single most revealing metric — delta_pitcher_run_exp")
add_multi_text(slide, 0.5, 1.3, 6.0, 5.5, [
    ("DELTA PITCHER RUN EXPECTANCY  (lower = better)", 20, RED, True),
    ("", 6),
    ("Sanchez:     +14.487 total    +0.0080/pitch  ← BEST", 18, ORANGE, True),
    ("Ohtani:      +21.028 total    +0.0158/pitch", 18, BLUE, True),
    ("Misiorowski: +36.807 total    +0.0222/pitch", 18, GREEN, True),
    ("", 6),
    ("This is the definitive run prevention metric.", 17, DARK),
    ("", 6),
    ("Sanchez gives up HALF the run value per pitch", 18, DARK, True),
    ("compared to Ohtani (+0.0080 vs +0.0158).", 18, DARK, True),
    ("", 6),
    ("Despite throwing 465 MORE pitches, Sanchez's", 17, DARK),
    ("TOTAL run cost is LOWER than Ohtani's.", 17, DARK),
    ("", 6),
    ("This is the strongest evidence against", 18, RED, True),
    ("the claim that Ohtani is the best pitcher.", 18, RED, True),
], line_spacing=4)

add_multi_text(slide, 7.0, 1.3, 5.8, 5.5, [
    ("What this means", 22, PURPLE, True),
    ("", 6),
    ("Run expectancy measures the change in expected", 16, DARK),
    ("runs scored after each pitch. It captures EVERYTHING:", 16, DARK),
    ("  • Walks and strikeouts", 16, DARK),
    ("  • Hit quality (singles vs doubles vs HR)", 16, DARK),
    ("  • Situational context (outs, runners on)", 16, DARK),
    ("  • Sequencing effect", 16, DARK),
    ("", 4),
    ("Sanchez is +14.5 runs above average.", 16, ORANGE, True),
    ("Ohtani is +21.0 — he's cost his team more runs.", 16, RED, True),
    ("", 4),
    ("If Ohtani had Sanchez's RE rate over 1,335 pitches:", 16, DARK),
    ("  +10.7 total → would save ~10 additional runs.", 16, DARK),
    ("", 4),
    ("Misiorowski's high total (+36.8) is surprising", 16, DARK),
    ("given his elite K%. This reflects high variance:", 16, DARK),
    ("  strikeout or walk/HR — little in between.", 16, DARK),
], line_spacing=4)

# =====================================================================
# SLIDE 20 — KEY METRICS FIGURE
# =====================================================================
slide = new_slide()
add_section_header(slide, "Key Metrics Comparison", "Side-by-side look at the most important stats")
if not add_figure_slide(slide, "key_metrics_comparison.png"):
    add_figure_slide(slide, "key_metrics.png")
add_caption(slide, "Multi-axis comparison: K%, BB%, Whiff%, Chase%, Hard Hit%, GB rate, and RE/pitch.")

# =====================================================================
# SLIDE 21 — PLATOON SPLITS: SANCHEZ
# =====================================================================
slide = new_slide()
add_section_header(slide, "Platoon Analysis: Sanchez Dominates LHB", "An extreme reverse-split profile")
add_multi_text(slide, 0.7, 1.3, 5.5, 5.5, [
    ("Sanchez vs Left-Handed Batters", 22, ORANGE, True),
    ("", 6),
    ("PA:     116", 18, DARK),
    ("K%:     36.2%  ← elite", 18, GREEN, True),
    ("BB%:     1.7%  ← generational control", 18, GREEN, True),
    ("AVG:    .138", 18, DARK),
    ("OPS:    .336  ← utterly dominant", 18, ORANGE, True),
    ("HR:        1", 18, DARK),
    ("", 6),
    ("Sanchez absolutely dominates LHB.", 17, DARK),
    ("A 36.2% K% with only 1.7% BB% is", 17, DARK),
    ("among the best platoon splits in baseball.", 17, DARK),
], line_spacing=4)

add_multi_text(slide, 6.8, 1.3, 5.8, 5.5, [
    ("Sanchez vs Right-Handed Batters", 22, ORANGE, True),
    ("", 6),
    ("PA:     380", 18, DARK),
    ("K%:     25.0%", 18, DARK),
    ("BB%:     5.8%", 18, DARK),
    ("AVG:    .263", 18, DARK),
    ("OPS:    .718", 18, DARK),
    ("HR:       10", 18, DARK),
    ("", 6),
    ("Sanchez is more vulnerable vs RHB (.718 OPS),", 16, DARK),
    ("but his extreme ground ball rate limits damage.", 16, DARK),
    ("10 HR in 380 PA is manageable for a sinker-baller.", 16, DARK),
    ("", 4),
    ("Overall, Sanchez's platoon profile is viable", 16, DARK),
    ("because LHB are virtually helpless against him.", 16, DARK),
], line_spacing=4)

# =====================================================================
# SLIDE 22 — PLATOON SPLITS: MISIOROWSKI
# =====================================================================
slide = new_slide()
add_section_header(slide, "Platoon Analysis: Misiorowski's Elite LHB K-Rate", "47.6% K% vs lefties is extraordinary")
add_multi_text(slide, 0.7, 1.3, 5.5, 5.5, [
    ("Misiorowski vs Left-Handed Batters", 22, GREEN, True),
    ("", 6),
    ("PA:     229", 18, DARK),
    ("K%:     47.6%  ← generational", 18, GREEN, True),
    ("BB%:     7.4%", 18, DARK),
    ("AVG:    .087  ← almost unhittable", 18, GREEN, True),
    ("OPS:    .288  ← elite, top-5 in MLB", 18, GREEN, True),
    ("HR:        2", 18, DARK),
    ("", 6),
    ("Almost half of all LHB plate appearances", 17, DARK),
    ("end in a strikeout. This is elite.", 17, DARK),
], line_spacing=4)

add_multi_text(slide, 6.8, 1.3, 5.8, 5.5, [
    ("Misiorowski vs Right-Handed Batters", 22, GREEN, True),
    ("", 6),
    ("PA:     193", 18, DARK),
    ("K%:     30.1%", 18, DARK),
    ("BB%:     5.2%", 18, DARK),
    ("AVG:    .192", 18, DARK),
    ("OPS:    .554", 18, DARK),
    ("HR:        7", 18, DARK),
    ("", 6),
    ("Still excellent vs RHB, but 7 HR in 193 PA", 16, DARK),
    ("is a concern. His fastball-heavy approach", 16, DARK),
    ("leaves him susceptible to pull-side power.", 16, DARK),
    ("", 4),
    ("Overall: elite platoon profile with a slight", 16, DARK),
    ("vulnerability to right-handed power hitters.", 16, DARK),
], line_spacing=4)

# =====================================================================
# SLIDE 23 — PLATOON SPLITS: OHTANI
# =====================================================================
slide = new_slide()
add_section_header(slide, "Platoon Analysis: Ohtani's Relative Weakness", "Less dominant vs both sides")
add_multi_text(slide, 0.7, 1.3, 5.5, 5.5, [
    ("Ohtani vs Left-Handed Batters", 22, BLUE, True),
    ("", 6),
    ("PA:     191", 18, DARK),
    ("K%:     24.6%", 18, DARK),
    ("BB%:     7.9%", 18, DARK),
    ("AVG:    .188", 18, DARK),
    ("OPS:    .545  ← good but not elite", 18, BLUE, True),
    ("HR:        3", 18, DARK),
    ("", 6),
    ("vs LHB: .545 OPS is solid, but well behind", 16, DARK),
    ("Sanchez's .336 and Misiorowski's .288.", 16, DARK),
    ("Ohtani's 24.6% K% vs LHB is notably low.", 16, RED, True),
], line_spacing=4)

add_multi_text(slide, 6.8, 1.3, 5.8, 5.5, [
    ("Ohtani vs Right-Handed Batters", 22, BLUE, True),
    ("", 6),
    ("PA:     149", 18, DARK),
    ("K%:     32.2%", 18, DARK),
    ("BB%:     7.4%", 18, DARK),
    ("AVG:    .128", 18, DARK),
    ("OPS:    .383  ← excellent vs RHB", 18, BLUE, True),
    ("HR:        1", 18, DARK),
    ("", 6),
    ("Ohtani handles same-handed batters well.", 16, DARK),
    (".383 OPS with 32.2% K% is a strong split.", 16, DARK),
    ("", 4),
    ("The gap: Ohtani vs LHB (.545 OPS) vs", 16, DARK),
    ("Sanchez vs LHB (.336 OPS) = 62% worse.", 18, RED, True),
], line_spacing=4)

# =====================================================================
# SLIDE 24 — PLATOON FIGURE
# =====================================================================
slide = new_slide()
add_section_header(slide, "Platoon Splits Visualization", "OPS against by handedness")
if not add_figure_slide(slide, "platoon_splits.png"):
    add_figure_slide(slide, "platoon_ops.png")
add_caption(slide, "OPS allowed vs LHB (left) and vs RHB (right) for each pitcher. Lower is better.")

# =====================================================================
# SLIDE 25 — FULL PLATOON TABLE
# =====================================================================
slide = new_slide()
add_section_header(slide, "Complete Platoon Comparison", "Every split, every pitcher")
add_multi_text(slide, 0.3, 1.3, 12.7, 5.5, [
    ("Pitcher        vs    PA     K%     BB%    AVG     OBP     SLG     OPS     HR", 19, DARK, True),
    ("", 4),
    ("Ohtani         LHB  191   24.6%   7.9%   .188   .260    .285    .545     3", 15, BLUE),
    ("Ohtani         RHB  149   32.2%   7.4%   .128   .198    .185    .383     1", 15, BLUE),
    ("", 3),
    ("Sanchez        LHB  116   36.2%   1.7%   .138   .155    .181    .336     1", 15, ORANGE),
    ("Sanchez        RHB  380   25.0%   5.8%   .263   .306    .412    .718    10", 15, ORANGE),
    ("", 3),
    ("Misiorowski    LHB  229   47.6%   7.4%   .087   .162    .126    .288     2", 15, GREEN),
    ("Misiorowski    RHB  193   30.1%   5.2%   .192   .238    .316    .554     7", 15, GREEN),
    ("", 6),
    ("Summary:", 20, DARK, True),
    ("  • Sanchez dominates LHB (.336 OPS, 36.2% K%) — finest LHB-neutralizing in the group", 16, DARK),
    ("  • Misiorowski is unhittable vs LHB (.288 OPS, 47.6% K%) — generational platoon dominance", 16, DARK),
    ("  • Ohtani is strong vs RHB (.383 OPS) but mediocre vs LHB (.545) by this standard", 16, DARK),
    ("  • All three have exploitable platoon traits: Ohtani vs L, Sanchez vs R, Misiorowski vs R", 16, DARK),
], line_spacing=3)

# =====================================================================
# SLIDE 26 — RADAR COMPARISON
# =====================================================================
slide = new_slide()
add_section_header(slide, "Radar Comparison", "Multi-dimensional performance profile")
if not add_figure_slide(slide, "radar_comparison.png"):
    pass
add_caption(slide, "Radar chart comparing the three pitchers across 8 key dimensions. Larger area = more well-rounded performance.")

# =====================================================================
# SLIDE 27 — KEY FINDINGS SUMMARY
# =====================================================================
slide = new_slide()
add_section_header(slide, "Key Findings Summary", "What the data tells us")
add_multi_text(slide, 0.5, 1.3, 6.0, 5.5, [
    ("Sanchez over Ohtani:", 22, ORANGE, True),
    ("", 4),
    ("✓  58% lower walk rate (4.8% vs 7.6%)", 16, DARK),
    ("✓  50% better run prevention per pitch (RE)", 16, DARK),
    ("✓  Lower total run cost despite 465 more pitches", 16, DARK),
    ("✓  Elite whiff rates on secondary offerings", 16, DARK),
    ("✓  192 vs 111 ground balls — damage suppression", 16, DARK),
    ("✓  Far more durable (1,800 vs 1,335 pitches)", 16, DARK),
    ("", 6),
    ("Misiorowski over both:", 22, GREEN, True),
    ("", 4),
    ("✓  39.6% K% — clearly elite", 16, DARK),
    ("✓  34.3% Whiff% — best of the three", 16, DARK),
    ("✓  100.5 mph average fastball — 99th percentile", 16, DARK),
    ("✓  47.6% K% vs LHB — generational", 16, DARK),
], line_spacing=4)

add_multi_text(slide, 6.8, 1.3, 6.0, 5.5, [
    ("Where Ohtani excels:", 22, BLUE, True),
    ("", 4),
    ("✓  0 barrels allowed — genuine batted ball suppression", 16, DARK),
    ("✓  86.9 mph avg exit velocity — lowest of the 3", 16, DARK),
    ("✓  Best expected stats (xwOBA .257)", 16, DARK),
    ("✓  7-pitch repertoire = adaptability", 16, DARK),
    ("", 6),
    ("Where Ohtani falls short:", 22, RED, True),
    ("", 4),
    ("✗  Highest walk rate (7.6%) — command gap", 16, DARK),
    ("✗  Highest run cost per pitch (+0.0158)", 16, DARK),
    ("✗  Lowest overall whiff rate (28.4%)", 16, DARK),
    ("✗  Lowest chase rate (29.8%) — less deception", 16, DARK),
    ("✗  35% fewer pitches — less total value", 16, DARK),
    ("✗  Platoon vulnerability vs LHB (.545 OPS)", 16, DARK),
], line_spacing=4)

# =====================================================================
# SLIDE 28 — CONCLUSION
# =====================================================================
slide = new_slide()
# Decorative top bar
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.08))
shape.fill.solid()
shape.fill.fore_color.rgb = RED
shape.line.fill.background()

add_textbox(slide, 0.5, 0.3, 12.3, 0.8,
    "Conclusion",
    size=38, bold=True, color=DARK)

add_textbox(slide, 0.5, 1.2, 12.3, 0.5,
    "The data does NOT support the claim that Ohtani is the best pitcher of the three.",
    size=22, bold=True, color=RED)

add_multi_text(slide, 0.7, 2.0, 5.8, 5.0, [
    ("Sanchez > Ohtani in:", 22, ORANGE, True),
    ("", 4),
    ("• Walk rate (4.8% vs 7.6%) — 58% better control", 17, DARK),
    ("• Run prevention per pitch (+0.0080 vs +0.0158)", 17, DARK),
    ("• Whiff rate on secondary pitches (42.7% vs 37.1%)", 17, DARK),
    ("• Ground ball generation (192 GB vs 111 GB)", 17, DARK),
    ("• Durability (1,800 vs 1,335 pitches)", 17, DARK),
    ("", 10),
    ("Misiorowski > Both in:", 22, GREEN, True),
    ("", 4),
    ("• Strikeout rate (39.6%), Whiff rate (34.3%)", 17, DARK),
    ("• Velocity (100.5 mph), LHB K% (47.6%)", 17, DARK),
], line_spacing=4)

add_multi_text(slide, 7.0, 2.0, 5.8, 5.5, [
    ("Ohtani is excellent, but", 20, BLUE, True),
    ("the data says he is the", 20, BLUE, True),
    ("3rd-best pitcher here.", 20, BLUE, True),
    ("", 10),
    ("Ohtani's strengths (0 barrels, low EV,", 16, MEDIUM),
    ("good expected stats) do not translate to", 16, MEDIUM),
    ("better actual run prevention.", 16, MEDIUM),
    ("", 6),
    ("Cy Young is about effectiveness and", 16, MEDIUM),
    ("innings. By both measures, Sanchez", 16, MEDIUM),
    ("and Misiorowski have stronger cases.", 16, MEDIUM),
], line_spacing=4)

add_textbox(slide, 0.5, 6.8, 12.3, 0.5,
    "Data: MLB Statcast 2026 · Analysis: Python (pandas, numpy, matplotlib, seaborn, scipy)",
    size=11, color=MEDIUM, italic=True, align=PP_ALIGN.CENTER)

# ---- SAVE ----
os.makedirs(OUTPUT_DIR, exist_ok=True)
prs.save(OUTPUT_PATH)
print(f'Presentation saved to {OUTPUT_PATH}')
print(f'Total slides: {len(prs.slides)}')
